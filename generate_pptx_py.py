"""
generate_pptx_py.py
===================
Pure-Python replacement for generate_pptx.js
Generates a 9-slide Sales ML Analytics PowerPoint with Ocean Dark theme.
Requires: python-pptx   (pip install python-pptx)

Usage (called directly from utils.py):
    from generate_pptx_py import build_presentation
    build_presentation(payload_dict, output_path)
"""

from __future__ import annotations
from pathlib import Path
from datetime import datetime
from typing import Any, Dict, List, Optional
import locale

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── PALETTE ────────────────────────────────────────────────────────────────────
def rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

P = {
    'bg':      '020B18',
    'bg2':     '03142E',
    'card':    '071828',
    'border':  '0D2D45',
    'accent':  '06B6D4',
    'sky':     '0EA5E9',
    'blue':    '38BDF8',
    'green':   '10B981',
    'amber':   'F59E0B',
    'red':     'EF4444',
    'purple':  '8B5CF6',
    'pink':    'EC4899',
    'white':   'FFFFFF',
    'text':    'E0F2FE',
    'muted':   '64748B',
    'chart_blues':  ['065A82','0D7FAD','0EA5E9','38BDF8','7DD3FC','BAE6FD'],
    'chart_greens': ['064E3B','065F46','047857','059669','10B981','34D399'],
    'chart_reds':   ['7F1D1D','991B1B','B91C1C','DC2626','EF4444','F87171'],
    'chart_multi':  ['06B6D4','10B981','F59E0B','8B5CF6','EC4899','EF4444','38BDF8','34D399'],
}

# ── SLIDE SIZE (16:9 widescreen) ───────────────────────────────────────────────
W = Inches(10)   # 10 in
H = Inches(5.625)  # 5.625 in

# ── FORMATTERS ─────────────────────────────────────────────────────────────────
def fmt_currency(v) -> str:
    try:
        v = float(v)
    except Exception:
        return "Rp 0"
    if v >= 1e12: return f"Rp {v/1e12:.1f}T"
    if v >= 1e9:  return f"Rp {v/1e9:.2f}M"
    if v >= 1e6:  return f"Rp {v/1e6:.1f}Jt"
    if v >= 1e3:  return f"Rp {v/1e3:.0f}K"
    return f"Rp {round(v)}"

def fmt_num(v) -> str:
    try:
        v = float(v)
    except Exception:
        return "0"
    if v >= 1e9: return f"{v/1e9:.1f}M"
    if v >= 1e6: return f"{v/1e6:.1f}Jt"
    if v >= 1e3: return f"{v/1e3:.0f}K"
    return str(round(v))

def safe(v, fallback=0):
    try:
        f = float(v)
        if f != f or f == float('inf') or f == float('-inf'):
            return fallback
        return f
    except Exception:
        return fallback

def trunc(s: str, n: int = 18) -> str:
    s = str(s) if s else ""
    return s[:n-1] + "…" if len(s) > n else s


# ── LOW-LEVEL HELPERS ──────────────────────────────────────────────────────────
def i(v): return Inches(v)   # shorthand

def add_rect(slide, x, y, w, h, fill_hex: str, line_hex: Optional[str] = None, transparency: int = 0):
    """Add a filled rectangle."""
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        i(x), i(y), i(w), i(h)
    )
    fill = shape.fill
    fill.solid()
    color = rgb(fill_hex)
    fill.fore_color.rgb = color

    if transparency > 0:
        # Set alpha via XML
        solidFill = shape.fill._xPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is None:
                srgbClr = solidFill.find(qn('a:sysClr'))
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha.set('val', str(int((100 - transparency) * 1000)))

    line = shape.line
    if line_hex:
        line.color.rgb = rgb(line_hex)
        line.width = Pt(0.75)
    else:
        line.fill.background()
    return shape


def add_text_box(slide, text: str, x, y, w, h,
                 font_size: int = 12,
                 bold: bool = False,
                 italic: bool = False,
                 color_hex: str = 'FFFFFF',
                 align=PP_ALIGN.LEFT,
                 font_name: str = 'Calibri',
                 wrap: bool = True):
    txBox = slide.shapes.add_textbox(i(x), i(y), i(w), i(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)
    run.font.name = font_name
    return txBox


def set_slide_bg(slide, prs, color_hex: str):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_hex)


def add_base(slide, prs):
    """Dark background + top accent line."""
    set_slide_bg(slide, prs, P['bg'])
    add_rect(slide, 0, 0, 10, 0.055, P['accent'])


def add_footer(slide, prs, page_num: int, total: int):
    FOOTER_Y = 5.3
    add_rect(slide, 0, FOOTER_Y, 10, 0.22, P['bg2'], P['border'])
    add_text_box(slide, "SalesML Analytics Pro  |  Confidential",
                 0.35, FOOTER_Y + 0.03, 6, 0.16,
                 font_size=8, color_hex=P['muted'])
    ts = datetime.now().strftime('%d %B %Y')
    add_text_box(slide, ts,
                 6.3, FOOTER_Y + 0.03, 2.4, 0.16,
                 font_size=8, color_hex=P['muted'], align=PP_ALIGN.RIGHT)
    add_text_box(slide, f"{page_num} / {total}",
                 8.8, FOOTER_Y + 0.03, 0.9, 0.16,
                 font_size=8, color_hex=P['accent'], align=PP_ALIGN.RIGHT)


def sec_label(slide, prs, text: str, x, y, w=4.5):
    """Section label with left accent bar."""
    add_rect(slide, x, y, 0.055, 0.28, P['accent'])
    add_text_box(slide, text, x + 0.1, y, w - 0.1, 0.28,
                 font_size=10, bold=True, color_hex=P['blue'])


def add_card(slide, prs, x, y, w, h):
    """Dark card with border."""
    add_rect(slide, x, y, w, h, P['card'], P['border'])


def kpi_block(slide, prs, x, y, w, h, icon: str, value: str, label: str,
              color_hex: str, delta: str = "", delta_up: bool = True):
    add_card(slide, prs, x, y, w, h)
    add_rect(slide, x, y, w, 0.04, color_hex)
    add_text_box(slide, icon,     x+0.15, y+0.1,  0.55, 0.45, font_size=18)
    add_text_box(slide, value,    x+0.15, y+0.5,  w-0.25, 0.4,
                 font_size=16, bold=True, color_hex=P['white'])
    add_text_box(slide, label.upper(), x+0.15, y+0.88, w-0.25, 0.22,
                 font_size=7, color_hex=P['blue'])
    if delta:
        add_text_box(slide, delta, x+0.15, y+1.08, w-0.25, 0.2,
                     font_size=9, bold=True,
                     color_hex=P['green'] if delta_up else P['red'])


def insight_row(slide, prs, text: str, x, y, w):
    add_card(slide, prs, x, y, w, 0.44)
    add_rect(slide, x, y, 0.055, 0.44, P['amber'])
    # Strip leading emoji/symbol
    import re
    clean = re.sub(r'^[^\w\s]+\s*', '', str(text))
    add_text_box(slide, clean, x+0.12, y+0.04, w-0.2, 0.36,
                 font_size=9, color_hex=P['text'], wrap=True)


def _apply_chart_theme(chart, chart_frame, values=None):
    """Apply dark theme + readable fonts to any chart. Shared by all chart types.
    chart = Chart object, chart_frame = GraphicFrame (parent shape)
    """
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    # Background — access via chart directly
    try:
        chart.chart_area.format.fill.solid()
        chart.chart_area.format.fill.fore_color.rgb = rgb(P['card'])
        chart.plot_area.format.fill.solid()
        chart.plot_area.format.fill.fore_color.rgb = rgb(P['card'])
    except Exception:
        try:
            # fallback via chart_frame if chart is actually GraphicFrame
            chart_frame.chart.chart_area.format.fill.solid()
            chart_frame.chart.chart_area.format.fill.fore_color.rgb = rgb(P['card'])
        except Exception:
            pass

    # Kill chart title
    try:
        chart.has_title = False
    except Exception:
        pass
    try:
        chart_elem = chart._element
        for title_elem in chart_elem.findall('.//' + qn('c:title')):
            title_elem.getparent().remove(title_elem)
    except Exception:
        pass

    # Axis tick labels — small, muted
    for axis_attr in ['category_axis', 'value_axis']:
        try:
            axis = getattr(chart, axis_attr)
            axis.tick_labels.font.color.rgb = rgb(P['muted'])
            axis.tick_labels.font.size = Pt(7)
        except Exception:
            pass

    # Gridlines — subtle
    try:
        chart.value_axis.major_gridlines.format.line.color.rgb = rgb(P['border'])
    except Exception:
        pass

    # Legend off
    try:
        chart.has_legend = False
    except Exception:
        pass

    # Data labels — disable entirely via XML
    try:
        chart_elem = chart._element
        for dlbls in chart_elem.iter(qn('c:dLbls')):
            for sv in dlbls.findall(qn('c:showVal')):
                sv.set('val', '0')
            for sc in dlbls.findall(qn('c:showCatName')):
                sc.set('val', '0')
            for sp in dlbls.findall(qn('c:showPercent')):
                sp.set('val', '0')
    except Exception:
        pass

    # Format value axis numbers (shorten big numbers like 200000000 → 200Jt)
    # Must inject numFmt directly into valAx XML
    if values:
        try:
            from pptx.oxml.ns import qn as _qn
            import lxml.etree as _etree
            max_val = max(abs(float(v)) for v in values if v is not None) if values else 1
            if max_val >= 1e9:
                fmt = '#,##0,,,"M"'
            elif max_val >= 1e6:
                fmt = '#,##0,,"Jt"'
            elif max_val >= 1e3:
                fmt = '#,##0,"K"'
            else:
                fmt = '#,##0.#'

            chart_elem = chart._element
            for valAx in chart_elem.iter(_qn('c:valAx')):
                # Remove existing numFmt if any
                for nf in valAx.findall(_qn('c:numFmt')):
                    valAx.remove(nf)
                # Create and insert numFmt at position 4 (after axPos)
                numFmt = _etree.Element(_qn('c:numFmt'))
                numFmt.set('formatCode', fmt)
                numFmt.set('sourceLinked', '0')
                valAx.insert(4, numFmt)
        except Exception:
            pass

    # Delete chart title: set autoTitleDeleted=1 and remove any c:title
    try:
        from pptx.oxml.ns import qn as _qn
        chart_elem = chart._element  # this is c:chartSpace
        c_chart = chart_elem.find(_qn('c:chart'))  # c:chart is child of chartSpace
        if c_chart is not None:
            atd = c_chart.find(_qn('c:autoTitleDeleted'))
            if atd is not None:
                atd.set('val', '1')
            else:
                import lxml.etree as _etree
                atd = _etree.SubElement(c_chart, _qn('c:autoTitleDeleted'))
                atd.set('val', '1')
                c_chart.insert(0, atd)
            for title in c_chart.findall(_qn('c:title')):
                c_chart.remove(title)
    except Exception:
        pass

    # Remove legend via XML from c:chart
    try:
        from pptx.oxml.ns import qn as _qn
        chart_elem = chart._element
        c_chart = chart_elem.find(_qn('c:chart'))
        if c_chart is not None:
            for leg in c_chart.findall(_qn('c:legend')):
                c_chart.remove(leg)
    except Exception:
        pass


def add_bar_chart(slide, labels: list, values: list,
                  x, y, w, h,
                  colors: list,
                  chart_type=XL_CHART_TYPE.BAR_CLUSTERED,
                  series_name: str = ""):
    """Add a bar or column chart."""
    chart_data = ChartData()
    chart_data.categories = [str(l) for l in labels]
    chart_data.add_series(series_name, [float(v) for v in values])

    chart_frame = slide.shapes.add_chart(
        chart_type, i(x), i(y), i(w), i(h), chart_data
    )
    chart = chart_frame.chart

    # Color each bar/column point
    plot = chart.plots[0]
    for idx, point in enumerate(plot.series[0].points):
        c = colors[idx % len(colors)]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = rgb(c)

    _apply_chart_theme(chart, chart_frame, values)

    return chart_frame


def add_line_chart(slide, labels: list, values: list, x, y, w, h, color_hex: str):
    chart_data = ChartData()
    chart_data.categories = [str(l) for l in labels]
    chart_data.add_series("", [float(v) for v in values])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, i(x), i(y), i(w), i(h), chart_data
    )
    chart = chart_frame.chart

    try:
        series = chart.plots[0].series[0]
        series.format.line.color.rgb = rgb(color_hex)
        series.format.line.width = Pt(2)
    except Exception:
        pass

    _apply_chart_theme(chart, chart_frame, values)
    return chart_frame


def add_doughnut_chart(slide, labels: list, values: list, x, y, w, h, colors: list):
    chart_data = ChartData()
    chart_data.categories = [str(l) for l in labels]
    chart_data.add_series("", [float(v) for v in values])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, i(x), i(y), i(w), i(h), chart_data
    )
    chart = chart_frame.chart

    try:
        for idx, point in enumerate(chart.plots[0].series[0].points):
            c = colors[idx % len(colors)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = rgb(c)
    except Exception:
        pass

    _apply_chart_theme(chart, chart_frame)

    # Re-enable legend for doughnut only
    try:
        chart.has_legend = True
        chart.legend.font.color.rgb = rgb(P['text'])
        chart.legend.font.size = Pt(8)
    except Exception:
        pass

    return chart_frame


# ── SLIDES ─────────────────────────────────────────────────────────────────────
TOTAL = 12


def slide_title(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, prs, P['bg'])

    # Decorative right panels
    add_rect(slide, 6.8, 0, 3.2, 5.625, P['sky'], P['bg'])
    add_rect(slide, 8.0, 0, 2.0, 5.625, P['accent'], P['bg'])

    # Top accent bar
    add_rect(slide, 0, 0, 10, 0.07, P['accent'])

    add_text_box(slide, "SalesML Analytics Pro",
                 0.5, 0.5, 6, 0.28, font_size=9, bold=True,
                 color_hex=P['accent'])
    add_text_box(slide, "Sales Performance\nReport",
                 0.5, 0.9, 6.1, 1.65, font_size=38, bold=True,
                 color_hex=P['white'])

    dr = data.get('date_range') or "Periode Analisis Data"
    add_text_box(slide, dr, 0.5, 2.65, 6, 0.3,
                 font_size=13, italic=True, color_hex=P['blue'])

    add_rect(slide, 0.5, 3.07, 2.2, 0.035, P['accent'])

    k = data.get('kpis') or {}
    pills = [
        ("Revenue",   fmt_currency(k.get('total_revenue', 0))),
        ("Transaksi", fmt_num(k.get('total_transactions', 0))),
        ("Produk",    fmt_num(k.get('unique_products', 0))),
    ]
    for idx, (label, val) in enumerate(pills):
        px = 0.5 + idx * 2.1
        add_card(slide, prs, px, 3.28, 1.9, 0.72)
        add_text_box(slide, val,   px+0.1, 3.34, 1.7, 0.3, font_size=13, bold=True, color_hex=P['blue'])
        add_text_box(slide, label.upper(), px+0.1, 3.64, 1.7, 0.22, font_size=7, color_hex=P['muted'])

    ts = datetime.now().strftime('%d %B %Y')
    add_text_box(slide, f"Generated: {ts}", 0.5, 5.05, 6, 0.22, font_size=8, color_hex=P['muted'])
    add_footer(slide, prs, 1, TOTAL)


def slide_agenda(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, prs)
    add_text_box(slide, "Agenda", 0.5, 0.22, 9, 0.55,
                 font_size=26, bold=True, color_hex=P['white'])

    items = [
        ("01", "Executive Summary",      "KPI utama, growth & key insights"),
        ("02", "Tren Revenue",            "Tren bulanan & MoM growth"),
        ("03", "Sales Performance",       "Top & bottom produk by revenue"),
        ("04", "Profitability Analysis",  "Margin & profit per produk"),
        ("05", "Customer & RFM",          "Segmentasi & top customer"),
        ("06", "Regional & Category",     "Wilayah, kategori & Pareto 80/20"),
        ("07", "Branch / Regional Deep",  "Revenue, growth & ranking per cabang"),
        ("08", "Channel Analysis",        "Performa & growth per channel penjualan"),
        ("09", "Sales Person",            "Ranking & kontribusi per sales person"),
    ]
    # 3-column layout for 9 items
    for i_idx, (num, title, sub) in enumerate(items):
        col = i_idx % 3
        row = i_idx // 3
        x = 0.28 + col * 3.18
        y = 1.0 + row * 1.3
        add_card(slide, prs, x, y, 3.05, 1.12)
        add_rect(slide, x, y, 0.055, 1.12, P['accent'])
        add_text_box(slide, num,   x+0.13, y+0.1,  0.5,  0.38, font_size=16, bold=True, color_hex=P['accent'])
        add_text_box(slide, title, x+0.68, y+0.1,  2.3,  0.38, font_size=11, bold=True, color_hex=P['white'])
        add_text_box(slide, sub,   x+0.68, y+0.52, 2.3,  0.42, font_size=8,  color_hex=P['muted'])

    add_footer(slide, prs, 2, TOTAL)


def slide_kpi(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, prs)
    add_text_box(slide, "Executive Summary", 0.5, 0.2, 9, 0.48,
                 font_size=24, bold=True, color_hex=P['white'])
    add_text_box(slide, "Ringkasan performa bisnis keseluruhan",
                 0.5, 0.65, 9, 0.25, font_size=11, color_hex=P['muted'])

    k = data.get('kpis') or {}
    g = data.get('growth') or {}
    margin = safe(k.get('avg_margin', 30), 30)
    profit = safe(k.get('total_revenue', 0)) * margin / 100
    mom = g.get('mom')
    yoy = g.get('yoy')
    total_rev = safe(k.get('total_revenue', 0))
    total_txn = max(safe(k.get('total_transactions', 1)), 1)
    aov = safe(k.get('avg_order_value', total_rev / total_txn))

    mom_str = (f"{'+' if mom and mom>=0 else ''}{mom:.1f}% vs bln lalu" if mom is not None else "")
    yoy_str = (f"{'+' if yoy and yoy>=0 else ''}{yoy:.1f}%" if yoy is not None else "N/A")
    mom_disp = (f"{'+' if mom>=0 else ''}{mom:.1f}%" if mom is not None else "N/A")

    cards = [
        ("💰", fmt_currency(total_rev), "Total Revenue",     P['accent'],  mom_str,  True),
        ("📦", fmt_num(k.get('total_quantity',0)), "Total Qty Terjual", P['green'], "", True),
        ("🧾", fmt_num(k.get('total_transactions',0)), "Total Transaksi", P['purple'], "", True),
        ("🛒", fmt_currency(aov),       "Avg Order Value",   P['amber'],   "", True),
        ("💎", fmt_currency(profit),    "Gross Profit",      P['green'],   "", True),
        ("📈", f"{margin:.1f}%",        "Profit Margin",     P['sky'],     "", True),
        ("📅", mom_disp,                "Growth MoM",        P['green'] if mom and mom>=0 else P['red'], "", True),
        ("🗓", yoy_str,                 "Growth YoY",        P['green'] if yoy and yoy>=0 else P['red'], "", True),
    ]
    ROW_Y = [1.0, 2.52]
    for idx, (icon, val, label, color, delta, dup) in enumerate(cards):
        row = idx // 4
        col = idx % 4
        kpi_block(slide, prs, 0.3 + col*2.38, ROW_Y[row], 2.22, 1.38,
                  icon, val, label, color, delta, dup)

    insights = (data.get('insights') or [])[:2]
    if insights:
        sec_label(slide, prs, "Key Insights", 0.3, 4.03, 9.4)
        for i_idx, ins in enumerate(insights):
            insight_row(slide, prs, ins, 0.3 + i_idx*4.87, 4.37, 4.65)

    add_footer(slide, prs, 3, TOTAL)


def slide_trend(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, prs)
    add_text_box(slide, "Tren Revenue", 0.5, 0.2, 9, 0.48,
                 font_size=24, bold=True, color_hex=P['white'])

    trend = data.get('monthly_trend') or []
    if not trend:
        add_text_box(slide, "Data tren tidak tersedia",
                     3, 2.5, 4, 0.5, font_size=14, color_hex=P['muted'],
                     align=PP_ALIGN.CENTER)
        add_footer(slide, prs, 4, TOTAL)
        return

    labels   = [t.get('month','') for t in trend]
    revenues = [round(safe(t.get('revenue', 0))) for t in trend]

    sec_label(slide, prs, "Revenue Bulanan", 0.3, 0.82, 9.4)
    try:
        add_bar_chart(slide, labels, revenues, 0.3, 1.08, 9.4, 2.2,
                      P['chart_blues'], XL_CHART_TYPE.COLUMN_CLUSTERED,
                      "")
    except Exception:
        pass

    mom_vals   = [round(safe(t.get('revenue_mom', 0), 0), 1) for t in trend[1:]]
    mom_labels = labels[1:]
    if len(mom_vals) > 1:
        sec_label(slide, prs, "MoM Growth (%)", 0.3, 3.4, 9.4)
        try:
            add_line_chart(slide, mom_labels, mom_vals, 0.3, 3.65, 9.4, 1.45, P['amber'])
        except Exception:
            pass

    add_footer(slide, prs, 4, TOTAL)


def slide_sales(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, prs)
    add_text_box(slide, "Sales Performance", 0.5, 0.2, 9, 0.48,
                 font_size=24, bold=True, color_hex=P['white'])
    add_text_box(slide, "Top & bottom produk berdasarkan revenue",
                 0.5, 0.65, 9, 0.25, font_size=11, color_hex=P['muted'])

    top_p = (data.get('top_products') or [])[:8]
    bot_p = (data.get('bottom_products') or [])[:8]

    if top_p:
        sec_label(slide, prs, "Top 8 Produk by Revenue", 0.3, 1.0, 4.6)
        try:
            add_bar_chart(slide,
                          [trunc(p.get('product',''), 20) for p in top_p],
                          [round(safe(p.get('revenue',0))) for p in top_p],
                          0.3, 1.28, 4.6, 3.75,
                          P['chart_blues'], XL_CHART_TYPE.BAR_CLUSTERED)
        except Exception:
            pass

    if bot_p:
        sec_label(slide, prs, "Bottom 8 — Perlu Perhatian", 5.1, 1.0, 4.6)
        try:
            add_bar_chart(slide,
                          [trunc(p.get('product',''), 20) for p in bot_p],
                          [round(safe(p.get('revenue',0))) for p in bot_p],
                          5.1, 1.28, 4.6, 3.75,
                          P['chart_reds'], XL_CHART_TYPE.BAR_CLUSTERED)
        except Exception:
            pass

    add_footer(slide, prs, 5, TOTAL)


def slide_profit(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, prs)
    add_text_box(slide, "Profitability Analysis", 0.5, 0.2, 9, 0.48,
                 font_size=24, bold=True, color_hex=P['white'])

    prods = data.get('profit_by_product') or []
    if not prods:
        add_text_box(slide, "Data profit tidak tersedia",
                     3, 2.5, 4, 0.5, font_size=14, color_hex=P['muted'],
                     align=PP_ALIGN.CENTER)
        add_footer(slide, prs, 6, TOTAL)
        return

    sorted_desc = sorted(prods, key=lambda p: safe(p.get('margin_pct',0)), reverse=True)
    sorted_asc  = sorted(prods, key=lambda p: safe(p.get('margin_pct',0)))
    top8 = sorted_desc[:8]
    bot8 = sorted_asc[:8]

    avg_m = sum(safe(p.get('margin_pct',0)) for p in prods) / len(prods)
    total_profit = sum(safe(p.get('profit',0)) for p in prods)
    total_rev    = sum(safe(p.get('revenue',0)) for p in prods)

    kpi_block(slide, prs, 0.3,      0.82, 3.0, 1.05, "💎", fmt_currency(total_profit), "Total Profit",   P['green'])
    kpi_block(slide, prs, 0.3+3.18, 0.82, 3.0, 1.05, "📈", f"{avg_m:.1f}%",            "Avg Margin",     P['sky'])
    kpi_block(slide, prs, 0.3+6.36, 0.82, 3.0, 1.05, "💰", fmt_currency(total_rev),    "Total Revenue",  P['accent'])

    sec_label(slide, prs, "Top 8 Margin %", 0.3, 1.98, 4.6)
    try:
        add_bar_chart(slide,
                      [trunc(p.get('product',''), 20) for p in top8],
                      [round(safe(p.get('margin_pct',0)), 1) for p in top8],
                      0.3, 2.26, 4.6, 2.85,
                      P['chart_greens'], XL_CHART_TYPE.BAR_CLUSTERED)
    except Exception:
        pass

    sec_label(slide, prs, "Margin Terendah (Waspada!)", 5.1, 1.98, 4.6)
    try:
        add_bar_chart(slide,
                      [trunc(p.get('product',''), 20) for p in bot8],
                      [round(safe(p.get('margin_pct',0)), 1) for p in bot8],
                      5.1, 2.26, 4.6, 2.85,
                      P['chart_reds'], XL_CHART_TYPE.BAR_CLUSTERED)
    except Exception:
        pass

    add_footer(slide, prs, 6, TOTAL)


def slide_customer(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, prs)
    add_text_box(slide, "Customer & RFM Analysis", 0.5, 0.2, 9, 0.48,
                 font_size=24, bold=True, color_hex=P['white'])

    segs  = data.get('rfm_segments') or {}
    top_c = (data.get('top_customers') or [])[:8]

    if segs:
        seg_keys = list(segs.keys())
        seg_vals = [segs[k] for k in seg_keys]
        sec_label(slide, prs, "Distribusi Segmen RFM", 0.3, 0.82, 4.6)
        try:
            add_doughnut_chart(slide, seg_keys, seg_vals,
                               0.3, 1.1, 4.6, 3.95,
                               [P['accent'],P['green'],P['amber'],P['purple'],P['red']])
        except Exception:
            pass

    if top_c:
        sec_label(slide, prs, "Top Customer by Revenue", 5.1, 0.82, 4.6)
        # Simple text table
        add_card(slide, prs, 5.1, 1.1, 4.6, 3.95)
        headers = ["#", "Customer", "Revenue", "Freq"]
        col_widths = [0.35, 1.8, 1.5, 0.9]
        col_x_start = 5.15
        row_h = 0.44
        # Header row
        cx = col_x_start
        for hdr, cw in zip(headers, col_widths):
            add_rect(slide, cx, 1.1, cw, row_h, '0A1E35')
            add_text_box(slide, hdr, cx+0.05, 1.15, cw-0.1, row_h-0.1,
                         font_size=9, bold=True, color_hex=P['white'])
            cx += cw
        # Data rows
        for r_idx, cust in enumerate(top_c):
            ry = 1.1 + (r_idx+1)*row_h
            if ry + row_h > 5.1:
                break
            cx = col_x_start
            row_bg = P['card'] if r_idx % 2 == 0 else P['bg2']
            for col_val, cw, col_color in zip(
                [str(r_idx+1),
                 trunc(str(cust.get('customer','')), 18),
                 fmt_currency(cust.get('monetary', cust.get('revenue',0))),
                 str(cust.get('frequency',''))],
                col_widths,
                [P['blue'], P['text'], P['green'], P['muted']]
            ):
                add_rect(slide, cx, ry, cw, row_h, row_bg)
                add_text_box(slide, col_val, cx+0.05, ry+0.05, cw-0.1, row_h-0.1,
                             font_size=9, color_hex=col_color)
                cx += cw

    add_footer(slide, prs, 7, TOTAL)


def slide_regional_category(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, prs)
    add_text_box(slide, "Regional & Category", 0.5, 0.2, 9, 0.48,
                 font_size=24, bold=True, color_hex=P['white'])

    regional = (data.get('regional') or [])[:8]
    cats     = (data.get('categories') or [])[:8]

    if regional:
        sec_label(slide, prs, "Revenue per Wilayah", 0.3, 0.82, 4.6)
        try:
            add_bar_chart(slide,
                          [trunc(r.get('region', r.get('name','')), 18) for r in regional],
                          [round(safe(r.get('revenue',0))) for r in regional],
                          0.3, 1.1, 4.6, 2.5,
                          P['chart_blues'], XL_CHART_TYPE.BAR_CLUSTERED)
        except Exception:
            pass
        # Mini pie via text summary
        try:
            add_bar_chart(slide,
                          [trunc(r.get('region', r.get('name','')), 14) for r in regional],
                          [round(safe(r.get('revenue',0))) for r in regional],
                          0.3, 3.72, 4.6, 1.45,
                          P['chart_multi'], XL_CHART_TYPE.PIE)
        except Exception:
            pass

    if cats:
        sec_label(slide, prs, "Revenue per Kategori", 5.1, 0.82, 4.6)
        try:
            add_bar_chart(slide,
                          [trunc(c.get('category', c.get('name','')), 18) for c in cats],
                          [round(safe(c.get('revenue',0))) for c in cats],
                          5.1, 1.1, 4.6, 2.5,
                          P['chart_multi'], XL_CHART_TYPE.COLUMN_CLUSTERED)
        except Exception:
            pass

    par = data.get('pareto') or {}
    if par.get('top_product_count'):
        add_card(slide, prs, 5.1, 3.72, 4.6, 1.45)
        add_rect(slide, 5.1, 3.72, 0.055, 1.45, P['amber'])
        add_text_box(slide, "Pareto 80/20", 5.22, 3.79, 4.35, 0.27,
                     font_size=10, bold=True, color_hex=P['amber'])
        add_text_box(slide,
                     f"{par['top_product_count']} produk ({round(safe(par.get('pct',0)))}%) hasilkan 80% revenue.",
                     5.22, 4.1, 4.35, 0.28, font_size=11, bold=True, color_hex=P['white'])
        add_text_box(slide, par.get('insight',''), 5.22, 4.4, 4.35, 0.66,
                     font_size=9, color_hex=P['text'], wrap=True)

    add_footer(slide, prs, 8, TOTAL)


def _ranking_table(slide, prs, items: list, x, y, w, h,
                   col_defs: list, row_limit: int = 8):
    """
    Generic ranking table with alternating rows.
    col_defs = [{'key': ..., 'label': ..., 'w': float, 'color': hex, 'fmt': callable}]
    """
    if not items:
        return
    items = items[:row_limit]
    row_h = min(0.42, (h - 0.44) / max(len(items), 1))

    # Header
    add_rect(slide, x, y, w, 0.36, P['bg2'])
    cx = x + 0.06
    for cd in col_defs:
        add_text_box(slide, cd['label'], cx, y + 0.07, cd['w'] - 0.06, 0.24,
                     font_size=8, bold=True, color_hex=P['blue'])
        cx += cd['w']

    # Data rows
    for r_idx, item in enumerate(items):
        ry = y + 0.36 + r_idx * row_h
        if ry + row_h > y + h:
            break
        row_bg = P['card'] if r_idx % 2 == 0 else P['bg2']
        add_rect(slide, x, ry, w, row_h, row_bg)
        # left accent on top 3
        if r_idx < 3:
            medal_colors = [P['amber'], P['muted'], P['accent']]
            add_rect(slide, x, ry, 0.04, row_h, medal_colors[r_idx])
        cx = x + 0.06
        for cd in col_defs:
            val = item.get(cd['key'], '')
            fmt_fn = cd.get('fmt')
            display = fmt_fn(val) if fmt_fn and val != '' else str(val)
            add_text_box(slide, display, cx, ry + 0.06,
                         cd['w'] - 0.06, row_h - 0.1,
                         font_size=8, color_hex=cd.get('color', P['text']))
            cx += cd['w']


def _growth_badge(slide, prs, x, y, w, h, label: str, val_now, val_prev,
                  color: str, icon: str = ""):
    """Mini card showing metric + growth vs previous period."""
    add_card(slide, prs, x, y, w, h)
    add_rect(slide, x, y, w, 0.035, color)
    if icon:
        add_text_box(slide, icon, x + 0.1, y + 0.07, 0.4, 0.35, font_size=14)
    val_now_f   = safe(val_now)
    val_prev_f  = safe(val_prev)
    growth = ((val_now_f - val_prev_f) / val_prev_f * 100) if val_prev_f else 0
    arrow  = "▲" if growth >= 0 else "▼"
    g_color = P['green'] if growth >= 0 else P['red']
    add_text_box(slide, fmt_currency(val_now_f), x + 0.1, y + 0.38, w - 0.18, 0.32,
                 font_size=13, bold=True, color_hex=P['white'])
    add_text_box(slide, label.upper(), x + 0.1, y + 0.68, w - 0.18, 0.18,
                 font_size=7, color_hex=P['blue'])
    add_text_box(slide, f"{arrow} {abs(growth):.1f}% vs prev",
                 x + 0.1, y + 0.85, w - 0.18, 0.18,
                 font_size=7, bold=True, color_hex=g_color)


# ── SLIDE 10 · BRANCH / REGIONAL DEEP DIVE ────────────────────────────────────
def slide_branch(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, prs)
    add_text_box(slide, "Branch / Regional Performance", 0.3, 0.18, 9.4, 0.48,
                 font_size=22, bold=True, color_hex=P['white'])
    add_text_box(slide, "Revenue, growth & ranking per cabang/wilayah",
                 0.3, 0.63, 9.4, 0.22, font_size=10, color_hex=P['muted'])

    branches = data.get('branches') or data.get('branch_data') or []
    # Fallback: derive from regional if no branch data
    if not branches:
        branches = [
            {'name': r.get('region', r.get('name', '')),
             'revenue': r.get('revenue', 0),
             'revenue_prev': r.get('revenue_prev', safe(r.get('revenue', 0)) * 0.9),
             'transactions': r.get('transactions', 0),
             'target': r.get('target', safe(r.get('revenue', 0)) * 1.1)}
            for r in (data.get('regional') or [])
        ]

    if not branches:
        add_text_box(slide, "Data branch tidak tersedia",
                     3.5, 2.5, 3, 0.4, font_size=12, color_hex=P['muted'], align=PP_ALIGN.CENTER)
        add_footer(slide, prs, 10, TOTAL)
        return

    top_branches = sorted(branches, key=lambda b: safe(b.get('revenue', 0)), reverse=True)[:8]

    # ── TOP KPI CARDS ──────────────────────────────────────────────────────────
    total_rev   = sum(safe(b.get('revenue', 0)) for b in branches)
    total_prev  = sum(safe(b.get('revenue_prev', 0)) for b in branches)
    total_txn   = sum(safe(b.get('transactions', 0)) for b in branches)
    total_tgt   = sum(safe(b.get('target', 0)) for b in branches)
    ach_pct     = (total_rev / total_tgt * 100) if total_tgt else 0
    n_branches  = len(branches)

    kpi_w = 9.4 / 4
    for idx, (icon, val, label, col) in enumerate([
        ("🏢", fmt_currency(total_rev),   "Total Revenue",    P['accent']),
        ("📊", f"{ach_pct:.1f}%",          "Target Achievement", P['green']),
        ("🔢", fmt_num(total_txn),         "Total Transaksi",  P['sky']),
        ("🏬", str(n_branches),            "Jumlah Branch",    P['purple']),
    ]):
        bx = 0.3 + idx * kpi_w
        kpi_block(slide, prs, bx, 0.9, kpi_w - 0.12, 1.05, icon, val, label, col)

    # ── LEFT: Bar chart revenue per branch ────────────────────────────────────
    sec_label(slide, prs, "Revenue per Branch", 0.3, 2.1, 4.8)
    try:
        add_bar_chart(slide,
                      [trunc(b.get('name', ''), 16) for b in top_branches],
                      [round(safe(b.get('revenue', 0))) for b in top_branches],
                      0.3, 2.38, 4.8, 2.9,
                      P['chart_blues'], XL_CHART_TYPE.BAR_CLUSTERED)
    except Exception:
        pass

    # ── RIGHT: Ranking table ────────────────────────────────────────────────
    sec_label(slide, prs, "Ranking & Growth", 5.3, 2.1, 4.4)
    add_card(slide, prs, 5.3, 2.38, 4.4, 2.9)

    col_defs = [
        {'key': 'rank',        'label': '#',        'w': 0.3,  'color': P['accent']},
        {'key': 'name',        'label': 'Branch',   'w': 1.45, 'color': P['text'],
         'fmt': lambda v: trunc(str(v), 14)},
        {'key': 'revenue',     'label': 'Revenue',  'w': 1.35, 'color': P['blue'],
         'fmt': fmt_currency},
        {'key': 'growth_pct',  'label': 'Growth',   'w': 0.75, 'color': P['green'],
         'fmt': lambda v: (f"▲{v:.1f}%" if safe(v) >= 0 else f"▼{abs(safe(v)):.1f}%")},
        {'key': 'ach_pct',     'label': 'Ach%',     'w': 0.55, 'color': P['amber'],
         'fmt': lambda v: f"{safe(v):.0f}%"},
    ]

    # Enrich with computed fields
    enriched = []
    for r_idx, b in enumerate(top_branches):
        rev   = safe(b.get('revenue', 0))
        prev  = safe(b.get('revenue_prev', rev * 0.9))
        tgt   = safe(b.get('target', rev * 1.1))
        growth = ((rev - prev) / prev * 100) if prev else 0
        ach    = (rev / tgt * 100) if tgt else 0
        enriched.append({**b, 'rank': r_idx + 1,
                         'growth_pct': growth, 'ach_pct': ach})

    _ranking_table(slide, prs, enriched, 5.3, 2.38, 4.4, 2.9, col_defs, row_limit=8)

    add_footer(slide, prs, 10, TOTAL)


# ── SLIDE 11 · CHANNEL ANALYSIS ───────────────────────────────────────────────
def slide_channel(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, prs)
    add_text_box(slide, "Channel Analysis", 0.3, 0.18, 9.4, 0.48,
                 font_size=22, bold=True, color_hex=P['white'])
    add_text_box(slide, "Performa & growth per channel penjualan",
                 0.3, 0.63, 9.4, 0.22, font_size=10, color_hex=P['muted'])

    channels = data.get('channels') or data.get('channel_data') or []
    if not channels:
        add_text_box(slide, "Data channel tidak tersedia",
                     3.5, 2.5, 3, 0.4, font_size=12, color_hex=P['muted'], align=PP_ALIGN.CENTER)
        add_footer(slide, prs, 11, TOTAL)
        return

    channels_sorted = sorted(channels, key=lambda c: safe(c.get('revenue', 0)), reverse=True)

    # ── KPI CARDS ─────────────────────────────────────────────────────────────
    total_rev  = sum(safe(c.get('revenue', 0)) for c in channels)
    best       = channels_sorted[0] if channels_sorted else {}
    best_name  = best.get('name', best.get('channel', '—'))
    best_rev   = safe(best.get('revenue', 0))
    best_pct   = (best_rev / total_rev * 100) if total_rev else 0
    n_ch       = len(channels)

    # Avg growth
    growths = []
    for c in channels:
        rev  = safe(c.get('revenue', 0))
        prev = safe(c.get('revenue_prev', 0))
        if prev:
            growths.append((rev - prev) / prev * 100)
    avg_growth = sum(growths) / len(growths) if growths else 0

    kpi_w = 9.4 / 4
    for idx, (icon, val, label, col) in enumerate([
        ("💰", fmt_currency(total_rev),     "Total Revenue",    P['accent']),
        ("🥇", trunc(best_name, 12),        "Top Channel",      P['amber']),
        ("📈", f"{avg_growth:+.1f}%",       "Avg Growth",       P['green'] if avg_growth >= 0 else P['red']),
        ("📡", str(n_ch),                   "Jumlah Channel",   P['purple']),
    ]):
        bx = 0.3 + idx * kpi_w
        kpi_block(slide, prs, bx, 0.9, kpi_w - 0.12, 1.05, icon, val, label, col)

    # ── LEFT: Donut share per channel ─────────────────────────────────────────
    sec_label(slide, prs, "Revenue Share per Channel", 0.3, 2.1, 4.65)
    ch_labels = [trunc(c.get('name', c.get('channel', '')), 14) for c in channels_sorted[:7]]
    ch_vals   = [round(safe(c.get('revenue', 0))) for c in channels_sorted[:7]]
    try:
        add_doughnut_chart(slide, ch_labels, ch_vals,
                           0.3, 2.38, 4.65, 2.9, P['chart_multi'])
    except Exception:
        pass

    # ── RIGHT: Table with growth & conversion ─────────────────────────────────
    sec_label(slide, prs, "Detail per Channel", 5.15, 2.1, 4.55)
    add_card(slide, prs, 5.15, 2.38, 4.55, 2.9)

    col_defs = [
        {'key': 'rank',       'label': '#',         'w': 0.28, 'color': P['accent']},
        {'key': 'name',       'label': 'Channel',   'w': 1.35, 'color': P['text'],
         'fmt': lambda v: trunc(str(v), 13)},
        {'key': 'revenue',    'label': 'Revenue',   'w': 1.25, 'color': P['blue'],
         'fmt': fmt_currency},
        {'key': 'share_pct',  'label': 'Share',     'w': 0.65, 'color': P['amber'],
         'fmt': lambda v: f"{safe(v):.1f}%"},
        {'key': 'growth_pct', 'label': 'Growth',    'w': 1.02, 'color': P['green'],
         'fmt': lambda v: (f"▲{safe(v):.1f}%" if safe(v) >= 0 else f"▼{abs(safe(v)):.1f}%")},
    ]

    enriched = []
    for r_idx, c in enumerate(channels_sorted):
        rev  = safe(c.get('revenue', 0))
        prev = safe(c.get('revenue_prev', 0))
        gr   = ((rev - prev) / prev * 100) if prev else 0
        sh   = (rev / total_rev * 100) if total_rev else 0
        n    = c.get('name', c.get('channel', f'Channel {r_idx+1}'))
        enriched.append({'rank': r_idx + 1, 'name': n,
                         'revenue': rev, 'share_pct': sh, 'growth_pct': gr})

    _ranking_table(slide, prs, enriched, 5.15, 2.38, 4.55, 2.9, col_defs, row_limit=8)

    add_footer(slide, prs, 11, TOTAL)


# ── SLIDE 12 · SALES PERSON ───────────────────────────────────────────────────
def slide_salesperson(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, prs)
    add_text_box(slide, "Sales Person Performance", 0.3, 0.18, 9.4, 0.48,
                 font_size=22, bold=True, color_hex=P['white'])
    add_text_box(slide, "Ranking, kontribusi & growth per sales person",
                 0.3, 0.63, 9.4, 0.22, font_size=10, color_hex=P['muted'])

    salespeople = data.get('salespeople') or data.get('sales_persons') or data.get('salesperson_data') or []
    if not salespeople:
        add_text_box(slide, "Data sales person tidak tersedia",
                     3.5, 2.5, 3, 0.4, font_size=12, color_hex=P['muted'], align=PP_ALIGN.CENTER)
        add_footer(slide, prs, 12, TOTAL)
        return

    sp_sorted = sorted(salespeople, key=lambda s: safe(s.get('revenue', 0)), reverse=True)

    # ── KPI CARDS ─────────────────────────────────────────────────────────────
    total_rev  = sum(safe(s.get('revenue', 0)) for s in salespeople)
    top_sp     = sp_sorted[0] if sp_sorted else {}
    top_name   = top_sp.get('name', top_sp.get('salesperson', '—'))
    n_sp       = len(salespeople)
    avg_rev    = total_rev / n_sp if n_sp else 0

    total_prev = sum(safe(s.get('revenue_prev', 0)) for s in salespeople)
    team_growth = ((total_rev - total_prev) / total_prev * 100) if total_prev else 0

    kpi_w = 9.4 / 4
    for idx, (icon, val, label, col) in enumerate([
        ("💰", fmt_currency(total_rev),          "Total Revenue Tim",  P['accent']),
        ("🏆", trunc(top_name, 13),              "Top Performer",      P['amber']),
        ("📊", fmt_currency(avg_rev),            "Avg Revenue/Sales",  P['sky']),
        ("📈", f"{team_growth:+.1f}%",           "Team Growth",        P['green'] if team_growth >= 0 else P['red']),
    ]):
        bx = 0.3 + idx * kpi_w
        kpi_block(slide, prs, bx, 0.9, kpi_w - 0.12, 1.05, icon, val, label, col)

    # ── LEFT: Bar chart top 8 sales person ────────────────────────────────────
    top8 = sp_sorted[:8]
    sec_label(slide, prs, "Top Sales Person by Revenue", 0.3, 2.1, 4.8)
    try:
        add_bar_chart(slide,
                      [trunc(s.get('name', s.get('salesperson', '')), 14) for s in top8],
                      [round(safe(s.get('revenue', 0))) for s in top8],
                      0.3, 2.38, 4.8, 2.9,
                      P['chart_blues'], XL_CHART_TYPE.BAR_CLUSTERED)
    except Exception:
        pass

    # ── RIGHT: Ranking table with growth & target achievement ─────────────────
    sec_label(slide, prs, "Detail & Growth", 5.3, 2.1, 4.4)
    add_card(slide, prs, 5.3, 2.38, 4.4, 2.9)

    col_defs = [
        {'key': 'rank',       'label': '#',        'w': 0.28, 'color': P['accent']},
        {'key': 'name',       'label': 'Sales',    'w': 1.4,  'color': P['text'],
         'fmt': lambda v: trunc(str(v), 13)},
        {'key': 'revenue',    'label': 'Revenue',  'w': 1.25, 'color': P['blue'],
         'fmt': fmt_currency},
        {'key': 'growth_pct', 'label': 'Growth',   'w': 0.75, 'color': P['green'],
         'fmt': lambda v: (f"▲{safe(v):.1f}%" if safe(v) >= 0 else f"▼{abs(safe(v)):.1f}%")},
        {'key': 'ach_pct',    'label': 'Ach%',     'w': 0.72, 'color': P['amber'],
         'fmt': lambda v: f"{safe(v):.0f}%"},
    ]

    enriched = []
    for r_idx, s in enumerate(top8):
        rev  = safe(s.get('revenue', 0))
        prev = safe(s.get('revenue_prev', 0))
        tgt  = safe(s.get('target', rev * 1.1))
        gr   = ((rev - prev) / prev * 100) if prev else 0
        ach  = (rev / tgt * 100) if tgt else 0
        nm   = s.get('name', s.get('salesperson', f'SP {r_idx+1}'))
        enriched.append({'rank': r_idx + 1, 'name': nm,
                         'revenue': rev, 'growth_pct': gr, 'ach_pct': ach})

    _ranking_table(slide, prs, enriched, 5.3, 2.38, 4.4, 2.9, col_defs, row_limit=8)

    add_footer(slide, prs, 12, TOTAL)



def slide_closing(prs, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs, P['bg'])
    add_rect(slide, 0, 5.555,  10, 0.07, P['accent'])
    add_rect(slide, 0, 0,     3.8, 5.625, P['sky'], P['bg'])

    add_text_box(slide, "⚡",             1.2, 1.3,  1.4, 1.4, font_size=52, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Terima Kasih",   4.0, 1.25, 5.6, 0.85, font_size=36, bold=True, color_hex=P['white'])
    add_text_box(slide, "Data-driven decisions,\nbetter business outcomes.",
                 4.0, 2.18, 5.6, 0.85, font_size=14, italic=True, color_hex=P['blue'])
    add_rect(slide, 4.0, 3.18, 2.2, 0.04, P['accent'])

    k = data.get('kpis') or {}
    add_text_box(slide, f"Total Revenue: {fmt_currency(k.get('total_revenue',0))}",
                 4.0, 3.35, 5.6, 0.28, font_size=12, color_hex=P['text'])
    add_text_box(slide,
                 f"Total Transaksi: {fmt_num(k.get('total_transactions',0))}  |  Periode: {data.get('date_range','—')}",
                 4.0, 3.65, 5.6, 0.28, font_size=11, color_hex=P['muted'])
    add_text_box(slide, "Generated by SalesML Analytics Pro",
                 4.0, 4.15, 5.6, 0.26, font_size=9, color_hex=P['muted'])
    add_footer(slide, prs, TOTAL, TOTAL)


# ── MAIN ENTRY POINT ───────────────────────────────────────────────────────────
def build_presentation(data: dict, output_path: str) -> str:
    """
    Build the 12-slide Ocean Dark Sales presentation.

    Parameters
    ----------
    data : dict
        Payload with keys:
        kpis, growth, monthly_trend, top_products, bottom_products,
        profit_by_product, rfm_segments, top_customers,
        regional, categories, pareto, insights, date_range,
        branches / branch_data, channels / channel_data,
        salespeople / sales_persons / salesperson_data
    output_path : str
        Path to save the .pptx file.

    Returns
    -------
    str
        Absolute path of the written file.
    """
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs, data)             # 1
    slide_agenda(prs, data)            # 2
    slide_kpi(prs, data)               # 3
    slide_trend(prs, data)             # 4
    slide_sales(prs, data)             # 5
    slide_profit(prs, data)            # 6
    slide_customer(prs, data)          # 7
    slide_regional_category(prs, data) # 8
    slide_branch(prs, data)            # 9  ← NEW
    slide_channel(prs, data)           # 10 ← NEW
    slide_salesperson(prs, data)       # 11 ← NEW
    slide_closing(prs, data)           # 12

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)
